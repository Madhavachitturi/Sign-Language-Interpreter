from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── colour palette ──────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0D, 0x1B, 0x2A)
MID_BLUE    = RGBColor(0x1B, 0x48, 0x86)
ACCENT      = RGBColor(0x00, 0xB4, 0xD8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF4, 0xF8)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)

# ── helpers ──────────────────────────────────────────────────────
def fill_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, font_size=16, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=6):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

# ════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(slide, DARK_BLUE)

add_rect(slide, 0, 0, 13.33, 0.08, ACCENT)           # top bar
add_rect(slide, 0, 7.42, 13.33, 0.08, ACCENT)         # bottom bar
add_rect(slide, 0.5, 1.2, 12.33, 0.06, MID_BLUE)      # divider under title

add_text_box(slide, "✋  Sign Language Interpreter",
             0.5, 0.25, 12.33, 1.1,
             font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "ASL Alphabet Recognition using MediaPipe & Random Forest",
             0.5, 1.4, 12.33, 0.7,
             font_size=20, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text_box(slide, "A Real-Time Computer Vision Pipeline",
             0.5, 2.1, 12.33, 0.5,
             font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# team box
add_rect(slide, 1.5, 2.9, 10.33, 4.1, MID_BLUE)
add_rect(slide, 1.5, 2.9, 10.33, 0.45, ACCENT)

add_text_box(slide, "TEAM MEMBERS",
             1.5, 2.9, 10.33, 0.45,
             font_size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

members = [
    ("CHITTURI MADHAVA BADARI NARAYANA", "22B21A4222"),
    ("POKURI NANDINI",                   "22B21A4203"),
    ("VELUGUBANTLA SIVARAM",             "22B21A4232"),
    ("CHINTAPALLI KIRAN KUMAR",          "22B21A4257"),
    ("BOKKA VAMSI KRISHNA REDDY",        "22B21A4249"),
]
for i, (name, roll) in enumerate(members):
    y = 3.5 + i * 0.65
    add_text_box(slide, f"▸  {name}", 2.0, y, 7.5, 0.5,
                 font_size=15, bold=True, color=WHITE)
    add_text_box(slide, roll, 9.5, y, 2.0, 0.5,
                 font_size=14, color=ACCENT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 – Problem Statement & Motivation
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(slide, DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 7.42, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 0.08, 0.08, 7.34, MID_BLUE)

add_text_box(slide, "Problem Statement & Motivation",
             0.3, 0.15, 12.5, 0.7,
             font_size=28, bold=True, color=ACCENT)
add_rect(slide, 0.3, 0.85, 12.5, 0.05, MID_BLUE)

points = [
    ("Communication Barrier",
     "Over 70 million people worldwide use sign language as their primary form of communication, yet most hearing people cannot understand it."),
    ("Manual Interpretation is Costly",
     "Professional sign language interpreters are expensive and not always available — a real-time automated system bridges this gap."),
    ("Computer Vision Opportunity",
     "Advances in hand-landmark detection (MediaPipe) and lightweight ML classifiers (Random Forest) make accurate, CPU-only real-time recognition feasible."),
    ("Our Goal",
     "Build an end-to-end pipeline that captures webcam frames, extracts hand landmarks, and classifies all 26 ASL alphabets in real time with no GPU required."),
]

for i, (heading, body) in enumerate(points):
    y = 1.05 + i * 1.5
    add_rect(slide, 0.4, y, 12.2, 1.3, MID_BLUE)
    add_rect(slide, 0.4, y, 0.35, 1.3, ACCENT)
    add_text_box(slide, heading, 0.9, y + 0.07, 11.5, 0.4,
                 font_size=16, bold=True, color=ACCENT)
    add_text_box(slide, body, 0.9, y + 0.48, 11.3, 0.75,
                 font_size=13, color=LIGHT_GRAY, wrap=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 – System Architecture
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(slide, DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 7.42, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 0.08, 0.08, 7.34, MID_BLUE)

add_text_box(slide, "System Architecture — 4-Script Pipeline",
             0.3, 0.15, 12.5, 0.7,
             font_size=28, bold=True, color=ACCENT)
add_rect(slide, 0.3, 0.85, 12.5, 0.05, MID_BLUE)

stages = [
    ("01", "collect_imgs.py",        "Webcam → ./data/<class>/",
     "Opens webcam, prompts per class (press Q), captures 100 JPG frames per letter (A–Z) into numbered subdirectories."),
    ("02", "create_dataset.py",      "Images → data.pickle",
     "Runs MediaPipe (max 1 hand, confidence 0.3), extracts 21 landmarks, normalises to bounding-box origin → 42-float feature vector."),
    ("03", "train_classifier.py",    "data.pickle → model.p",
     "Loads pickle, splits 80/20, fits RandomForestClassifier (default params), prints accuracy, serialises model."),
    ("04", "inference_classifier.py","Webcam → Live prediction",
     "Loads model, loops over frames, extracts landmarks identically to Step 2, predicts letter, overlays bounding box + label."),
]

for i, (num, script, io, desc) in enumerate(stages):
    x = 0.35 + i * 3.24
    add_rect(slide, x, 1.05, 3.0, 5.9, MID_BLUE)
    add_rect(slide, x, 1.05, 3.0, 0.5, ACCENT)
    add_text_box(slide, num, x + 0.08, 1.07, 0.5, 0.4,
                 font_size=20, bold=True, color=DARK_BLUE)
    add_text_box(slide, script, x + 0.55, 1.07, 2.4, 0.42,
                 font_size=13, bold=True, color=DARK_BLUE)
    add_text_box(slide, io, x + 0.1, 1.65, 2.8, 0.55,
                 font_size=12, bold=True, color=ACCENT)
    add_text_box(slide, desc, x + 0.1, 2.3, 2.8, 4.5,
                 font_size=11, color=LIGHT_GRAY, wrap=True)
    if i < 3:
        add_text_box(slide, "→", x + 3.0, 3.5, 0.3, 0.5,
                     font_size=22, bold=True, color=ACCENT)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 – Feature Extraction & ML Model
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(slide, DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 7.42, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 0.08, 0.08, 7.34, MID_BLUE)

add_text_box(slide, "Feature Extraction & ML Model",
             0.3, 0.15, 12.5, 0.7,
             font_size=28, bold=True, color=ACCENT)
add_rect(slide, 0.3, 0.85, 12.5, 0.05, MID_BLUE)

# left column
add_rect(slide, 0.4, 1.0, 6.0, 5.9, MID_BLUE)
add_rect(slide, 0.4, 1.0, 6.0, 0.45, ACCENT)
add_text_box(slide, "MediaPipe Hand Landmark Extraction",
             0.5, 1.02, 5.8, 0.4,
             font_size=15, bold=True, color=DARK_BLUE)

left_points = [
    "• 21 key-points detected per hand (wrist + 4 per finger)",
    "• Each landmark has normalised (x, y) coordinates in frame space",
    "• Coordinates shifted by min(x) and min(y) → translation-invariant",
    "• Final feature vector: 42 floats  [x₀−xₘᵢₙ, y₀−yₘᵢₙ, ..., x₂₀−xₘᵢₙ, y₂₀−yₘᵢₙ]",
    "• Static image mode used during training (max 1 hand)",
    "• Same normalisation applied at inference — mismatch causes wrong predictions",
]
y = 1.55
for pt in left_points:
    add_text_box(slide, pt, 0.55, y, 5.7, 0.55,
                 font_size=12, color=LIGHT_GRAY, wrap=True)
    y += 0.75

# right column
add_rect(slide, 6.9, 1.0, 6.0, 5.9, MID_BLUE)
add_rect(slide, 6.9, 1.0, 6.0, 0.45, ACCENT)
add_text_box(slide, "Random Forest Classifier",
             7.0, 1.02, 5.8, 0.4,
             font_size=15, bold=True, color=DARK_BLUE)

right_points = [
    ("Why Random Forest?",
     "Handles non-linear feature relationships well; robust to noise in hand-landmark coordinates."),
    ("Training Split",
     "80% training / 20% test, stratified shuffle — ensures all 26 classes represented equally in both splits."),
    ("Accuracy",
     "100% on the held-out test set with the collected dataset."),
    ("Serialisation",
     "Model pickled to model.p via sklearn 1.2.0. Must be retrained if sklearn version changes — version mismatch crashes inference with a ValueError."),
    ("Inference Speed",
     "Runs in real time on CPU — no GPU required. Prediction latency is negligible compared to webcam frame rate."),
]
y = 1.55
for heading, body in right_points:
    add_text_box(slide, heading, 7.05, y, 5.7, 0.32,
                 font_size=13, bold=True, color=ACCENT)
    add_text_box(slide, body, 7.05, y + 0.32, 5.7, 0.55,
                 font_size=11, color=LIGHT_GRAY, wrap=True)
    y += 1.0

# ════════════════════════════════════════════════════════════════
# SLIDE 5 – Dataset
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(slide, DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 7.42, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 0.08, 0.08, 7.34, MID_BLUE)

add_text_box(slide, "Dataset",
             0.3, 0.15, 12.5, 0.7,
             font_size=28, bold=True, color=ACCENT)
add_rect(slide, 0.3, 0.85, 12.5, 0.05, MID_BLUE)

stats = [
    ("26",    "Classes\n(A – Z)"),
    ("100",   "Images\nper Class"),
    ("2,600", "Total\nFrames"),
    ("42",    "Features\nper Sample"),
]
for i, (num, label) in enumerate(stats):
    x = 0.5 + i * 3.2
    add_rect(slide, x, 1.0, 2.8, 1.8, MID_BLUE)
    add_rect(slide, x, 1.0, 2.8, 0.06, ACCENT)
    add_text_box(slide, num, x, 1.15, 2.8, 0.9,
                 font_size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, x, 1.95, 2.8, 0.7,
                 font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_rect(slide, 0.4, 3.05, 12.2, 4.0, MID_BLUE)
add_text_box(slide, "Collection & Processing Details",
             0.6, 3.1, 11.8, 0.45,
             font_size=16, bold=True, color=ACCENT)

details = [
    "▸  Captured live via webcam using OpenCV — no external dataset required.",
    "▸  100 JPEG frames per class stored in ./data/<class_index>/ (0 = A … 25 = Z).",
    "▸  MediaPipe processes each image in static_image_mode (single hand, confidence ≥ 0.3).",
    "▸  Images where no hand is detected are silently skipped — class imbalance possible if hand is consistently missed.",
    "▸  Landmarks normalised to the hand bounding-box origin before pickling — makes features scale/position invariant.",
    "▸  dataset_size is hardcoded to 100 in collect_imgs.py — change before re-collecting to increase dataset size.",
]
y = 3.6
for d in details:
    add_text_box(slide, d, 0.6, y, 12.0, 0.5,
                 font_size=13, color=LIGHT_GRAY, wrap=True)
    y += 0.55

# ════════════════════════════════════════════════════════════════
# SLIDE 6 – Tech Stack
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(slide, DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 7.42, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 0.08, 0.08, 7.34, MID_BLUE)

add_text_box(slide, "Tech Stack",
             0.3, 0.15, 12.5, 0.7,
             font_size=28, bold=True, color=ACCENT)
add_rect(slide, 0.3, 0.85, 12.5, 0.05, MID_BLUE)

techs = [
    ("Python 3.10", "Core language. Chosen for its rich ML/CV ecosystem and compatibility with all libraries used."),
    ("OpenCV 4.7",  "Webcam capture, frame reading, image I/O, bounding-box and text overlay rendering."),
    ("MediaPipe",   "Google's hand-landmark detection model. Outputs 21 3-D keypoints; only x, y used here."),
    ("scikit-learn 1.2.0", "RandomForestClassifier training, evaluation, and pickle serialisation. Version is pinned — must match model.p."),
    ("NumPy",       "Feature vector construction (np.asarray) and array manipulation for model input."),
    ("python-pptx", "Used to generate this presentation programmatically from the project codebase."),
]

for i, (tech, desc) in enumerate(techs):
    row, col = divmod(i, 2)
    x = 0.4 + col * 6.5
    y = 1.1 + row * 2.1
    add_rect(slide, x, y, 6.1, 1.85, MID_BLUE)
    add_rect(slide, x, y, 0.25, 1.85, ACCENT)
    add_text_box(slide, tech, x + 0.35, y + 0.15, 5.6, 0.45,
                 font_size=16, bold=True, color=ACCENT)
    add_text_box(slide, desc, x + 0.35, y + 0.6, 5.6, 1.1,
                 font_size=12, color=LIGHT_GRAY, wrap=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 – Known Issues & Future Improvements
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(slide, DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 7.42, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 0.08, 0.08, 7.34, MID_BLUE)

add_text_box(slide, "Known Issues & Future Improvements",
             0.3, 0.15, 12.5, 0.7,
             font_size=28, bold=True, color=ACCENT)
add_rect(slide, 0.3, 0.85, 12.5, 0.05, MID_BLUE)

add_rect(slide, 0.4, 1.0, 5.9, 5.9, MID_BLUE)
add_rect(slide, 0.4, 1.0, 5.9, 0.45, RGBColor(0xC0, 0x39, 0x2B))
add_text_box(slide, "⚠  Known Constraints",
             0.5, 1.02, 5.7, 0.4,
             font_size=15, bold=True, color=WHITE)

issues = [
    ("sklearn Version Lock",
     "model.p must be retrained with the same sklearn used at runtime. Version mismatch raises a hard ValueError at startup."),
    ("Single Hand Only",
     "Training uses max_num_hands=1 but inference defaults to 2. If two hands appear, feature vector doubles to 84 floats and prediction crashes."),
    ("Static Gestures Only",
     "Dynamic letters (J, Z) involve motion and cannot be captured by a single-frame landmark snapshot."),
    ("Fixed Dataset Size",
     "dataset_size=100 is hardcoded; small datasets may reduce generalisation across different hand sizes/lighting."),
]
y = 1.55
for heading, body in issues:
    add_text_box(slide, f"• {heading}", 0.55, y, 5.6, 0.35,
                 font_size=13, bold=True, color=RGBColor(0xFF, 0x6B, 0x6B))
    add_text_box(slide, body, 0.55, y + 0.35, 5.6, 0.6,
                 font_size=11, color=LIGHT_GRAY, wrap=True)
    y += 1.1

add_rect(slide, 7.0, 1.0, 5.9, 5.9, MID_BLUE)
add_rect(slide, 7.0, 1.0, 5.9, 0.45, RGBColor(0x1A, 0x73, 0x4A))
add_text_box(slide, "🚀  Future Improvements",
             7.1, 1.02, 5.7, 0.4,
             font_size=15, bold=True, color=WHITE)

futures = [
    ("Word & Sentence Detection",
     "Buffer consecutive letter predictions to form words; add space/backspace gestures."),
    ("Deep Learning Model",
     "Replace Random Forest with a CNN or LSTM to handle dynamic gestures (J, Z) and improve accuracy under varied lighting."),
    ("Text-to-Speech Output",
     "Pipe predicted text to a TTS engine (pyttsx3 / gTTS) for audio output."),
    ("Two-Hand Support",
     "Extend feature extraction to handle two-hand signs and fix the inference multi-hand bug."),
    ("Mobile / Web Deployment",
     "Export model to ONNX or TensorFlow Lite and deploy via a browser or Android app."),
]
y = 1.55
for heading, body in futures:
    add_text_box(slide, f"• {heading}", 7.15, y, 5.6, 0.35,
                 font_size=13, bold=True, color=RGBColor(0x2E, 0xCC, 0x71))
    add_text_box(slide, body, 7.15, y + 0.35, 5.6, 0.6,
                 font_size=11, color=LIGHT_GRAY, wrap=True)
    y += 1.05

# ════════════════════════════════════════════════════════════════
# SLIDE 8 – Thank You
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
fill_bg(slide, DARK_BLUE)
add_rect(slide, 0, 0, 13.33, 0.08, ACCENT)
add_rect(slide, 0, 7.42, 13.33, 0.08, ACCENT)
add_rect(slide, 3.5, 3.45, 6.33, 0.06, ACCENT)

add_text_box(slide, "Thank You",
             0, 1.5, 13.33, 1.3,
             font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "✋  Sign Language Interpreter — ASL A–Z Real-Time Recognition",
             0.5, 2.8, 12.33, 0.6,
             font_size=17, italic=True, color=ACCENT, align=PP_ALIGN.CENTER)

add_text_box(slide, "Team Members",
             0, 3.65, 13.33, 0.45,
             font_size=14, bold=True, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

names_line = "  ·  ".join(n for n, _ in members)
add_text_box(slide, names_line,
             0.3, 4.15, 12.73, 0.55,
             font_size=12, color=ACCENT, align=PP_ALIGN.CENTER)

add_text_box(slide, "Built with Python · OpenCV · MediaPipe · scikit-learn",
             0, 6.6, 13.33, 0.5,
             font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ── save ────────────────────────────────────────────────────────
out = "Sign_Language_Interpreter.pptx"
prs.save(out)
print(f"Saved → {out}")
